"""V4.pptx に赤字注記を追加 → V5.pptx として保存
（修正案として user に確認してもらう用、直接書き換えはしない）
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V4.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V5.pptx"

RED = RGBColor(0xC0, 0x00, 0x00)
RED_BG = RGBColor(0xFF, 0xEE, 0xEE)


def add_red_note(slide, x, y, w, h, text, size=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = 'Noto Sans JP'
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = RED
    tb.fill.solid()
    tb.fill.fore_color.rgb = RED_BG
    tb.line.color.rgb = RED
    tb.line.width = Pt(1.0)
    return tb


def main():
    prs = Presentation(str(SRC))
    slides = list(prs.slides)
    print(f"loaded V4: {len(slides)} slides")

    # Slide 6 (index 5): 「PB 空白カテゴリ仮説」
    add_red_note(
        slides[5], Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45),
        "⚠ FACT CHECK：ベース機種候補 トミー精工 FLS-1000 が既に「女性ユーザー着目・片手開閉アシスト・ローテーブル79cm設計」を NB として訴求済。仮説 A は『空白』ではなく『NB が占有』。真の差別化は FLS-1000 の弱点（バリデーション・データ可搬・GMP）補強へ再探索。"
    )

    # Slide 7 (index 6): 空白地帯マップ
    add_red_note(
        slides[6], Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45),
        "⚠ 「人間工学×中価格」の空白は FLS-1000 (¥1,080,000) が既に占有。真の空白地帯は「人間工学 + バリデーション/GMP」「人間工学 + IoT/データ連動」「人間工学 + サブスク」など弱点補強の組合せ軸で再描画必要。"
    )

    # Slide 9 (index 8): 「スペックちょい替え」woman-friendly 100L
    add_red_note(
        slides[8], Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.55),
        "⚠ PB 仮称「woman-friendly 100L」は FLS-1000 の design_concept そのもの（features にも「女性ユーザー」「片手開閉アシスト」明記）。PB として独自性を出すには NB 軸の継承＋以下いずれか必要：① バリデーション/IQ-OQ 強化, ② USB/LAN データ可搬, ③ 多言語 UI, ④ PB 価格圧縮 -15〜20%, ⑤ サブスク同梱。"
    )

    # Slide 10 (index 9): 想定スペック・訴求・価格
    add_red_note(
        slides[9], Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45),
        "⚠ 想定スペックの「電動アシスト上蓋・低位置投入口」は FLS-1000 が既に satisfy。キャッチ「重い蓋から、解放を。」「やさしいオートクレーブ」も NB 訴求の言い換え。PB 独自性を出すための新軸（バリデーション/データ/価格/サブスク）を再選定し、キャッチも刷新が必要。"
    )

    prs.save(str(DST))
    print(f"saved V5: {DST}")
    print("\n赤字注記を追加したスライド:")
    print("  Slide 6  PB 空白カテゴリ仮説")
    print("  Slide 7  ポジショニング・マトリクス")
    print("  Slide 9  スペックちょい替え")
    print("  Slide 10 想定スペック・訴求・価格")


if __name__ == "__main__":
    main()
