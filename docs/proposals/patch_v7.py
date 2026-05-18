"""V6 → V7: POS/SNS サンプル明示 + データ前提スライド追加

目的: 顧客（アズワン様）に渡す前のハルシネーション抑止強化。
- 既存スライドの POS/SNS 数値・ターゲット仮説に「サンプル」注記を併記
- 新規スライド「本提案書のデータ前提と PoC 移行時の置換ポリシー」を末尾に追加
- スクレイピング 101 件は実 EC データであることを明示し、線引きを明確化
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V6.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V7.pptx"


# patch_v4.py と同じ色・フォント
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_LIGHT = RGBColor(0x3A, 0x5A, 0x85)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)
ACCENT = RGBColor(0xC4, 0x76, 0x3A)
ACCENT_BG = RGBColor(0xFD, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x18, 0x27)
FONT = 'Noto Sans JP'

TITLE_PREFIX = '勘と手作業から脱却する、次世代PB開発の「設計図」'
FOOTER_TXT = 'アズワン株式会社様向け｜オートクレーブPB戦略提案｜2026年5月'


# Find/Replace（V6 の実テキストにマッチするもののみ）
# - Slide 7 のポジショニング section に「サンプル」注記
# - Slide 13 のリードタイム圧縮箇所に「PoC 想定サンプル」注記
# - Slide 6 / 9 / 10 の PB 仮説・ターゲット記述に「サンプル」併記
REPLACEMENTS = [
    # Slide 7 系：自社 POS × 外部スクレイピング
    ("自社 POS データ（内部）とスクレイピングデータ",
     "自社 POS データ（内部・実 PoC データ／本提案書はサンプル）とスクレイピングデータ"),

    # Slide 13 系：リードタイム圧縮の総括コメント
    ("浮いた時間を【単なる時短ではなく、マーケターの高度な戦略的意思決定に投資】可能。",
     "浮いた時間を【単なる時短ではなく、マーケターの高度な戦略的意思決定に投資】可能。\n※ 本提案書の数値は PoC 想定サンプル、実投入後は貴社実データで再生成。"),

    # Slide 6：PB 空白カテゴリ仮説（3 案）に「サンプル」を併記
    ("100L × GMP/IQ-OQ + データ可搬",
     "100L × GMP/IQ-OQ + データ可搬（※ サンプル仮説）"),
    ("100L × 多言語UI + クラウド通知 + サブスク",
     "100L × 多言語UI + クラウド通知 + サブスク（※ サンプル仮説）"),
    ("100L × 横置き省設置 + 消耗品同梱",
     "100L × 横置き省設置 + 消耗品同梱（※ サンプル仮説）"),

    # Slide 9 / 10 ターゲット記述（T-A / T-C）に「サンプル」併記
    ("T-A：製薬・CDMO 試作ラボ QA 担当（40代）／T-C：バイオベンチャー研究 PI（再生医療・治験原料）",
     "T-A：製薬・CDMO 試作ラボ QA 担当（40代）／T-C：バイオベンチャー研究 PI（再生医療・治験原料）　※ ターゲット像はサンプル、実 PoC では貴社実 POS で再定義"),

    # Slide 10 想定価格（NB 比 -5〜10%）にサンプル併記
    ("・想定価格：NB 比 -5〜10%（IQ-OQ込み総コストで -25%）",
     "・想定価格：NB 比 -5〜10%（IQ-OQ込み総コストで -25%）（※ サンプル試算）"),
]


# ---- ヘルパー（patch_v4.py から流用）----
def _set_font(run, size=14, bold=False, color=BLACK, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, border=None,
             border_w=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(border_w)
    else:
        tb.line.fill.background()
    return tb


def add_lines_box(slide, x, y, w, h, lines, fill=None, border=None, border_w=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln.get('text', '')
        _set_font(run,
                  size=ln.get('size', 12),
                  bold=ln.get('bold', False),
                  color=ln.get('color', BLACK))
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(border_w)
    else:
        tb.line.fill.background()
    return tb


def add_title(slide, title):
    add_text(slide, Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.4),
             TITLE_PREFIX, size=11, color=GRAY)
    add_text(slide, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.55),
             title, size=22, bold=True, color=NAVY)


def add_page_meta(slide, page_no, total=21):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8.0), Inches(0.3),
             FOOTER_TXT, size=9, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
             f'{page_no} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def _replace_in_runs(paragraph, old, new):
    full_text = "".join(run.text for run in paragraph.runs)
    if old not in full_text:
        return False
    new_full = full_text.replace(old, new)
    if not paragraph.runs:
        return False
    first = paragraph.runs[0]
    first.text = new_full
    for run in paragraph.runs[1:]:
        run.text = ""
    return True


def _walk_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        if hasattr(shape, 'shapes'):
            for s in shape.shapes:
                if hasattr(s, 'text_frame') and s.has_text_frame:
                    yield s.text_frame


# ---- 新規スライド: データ前提と置換ポリシー ----
def slide_21_data_premise(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_meta(s, 21, total=21)
    add_title(s, '本提案書のデータ前提と PoC 移行時の置換ポリシー')

    # 左：本提案書での扱い（サンプルデータ）
    add_text(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.4),
             '【本提案書での扱い】サンプルデータ', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(0.4), Inches(1.95), Inches(6.2), Inches(3.6),
                  fill=GRAY_BG, border=GRAY, border_w=0.75, lines=[
        {'text': '◆ POS 数値（PoC 想定例）', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・月間 35 台 / 100L 引き合い 1.4 倍', 'size': 11},
        {'text': '・価格帯 ¥40-90 万 60%', 'size': 11},
        {'text': '', 'size': 4},
        {'text': '◆ リピート率（PoC 想定例）', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・大学 30% / 製薬 50%', 'size': 11},
        {'text': '', 'size': 4},
        {'text': '◆ SNS の声（PoC 想定例）', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・女性研究者発信 60%（「蓋が重い」 等）', 'size': 11},
        {'text': '', 'size': 6},
        {'text': '→ これらは PoC 想定例として記載。', 'size': 11, 'color': ACCENT, 'bold': True},
        {'text': '　 実 PoC では貴社実 POS / SNS データに置換します。', 'size': 11, 'color': ACCENT, 'bold': True},
    ])

    # 右：実 PoC では（置換ポリシー）
    add_text(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.4),
             '【実 PoC では】実データに置換するもの', size=14, bold=True, color=ACCENT)
    add_lines_box(s, Inches(6.8), Inches(1.95), Inches(6.2), Inches(3.6),
                  fill=ACCENT_BG, border=ACCENT, border_w=0.75, lines=[
        {'text': '◆ アズワン社内 POS', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・購買履歴・属性・行動データ', 'size': 11},
        {'text': '', 'size': 4},
        {'text': '◆ 既存 PB ラインアップ', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・実カタログ・商品マスタ', 'size': 11},
        {'text': '', 'size': 4},
        {'text': '◆ SNS 観測', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・実クエリ × 実発信内容', 'size': 11},
        {'text': '', 'size': 4},
        {'text': '◆ 業界レポート', 'size': 12, 'bold': True, 'color': NAVY},
        {'text': '・矢野経済研・富士経済など（有償）', 'size': 11},
        {'text': '', 'size': 6},
        {'text': '→ PoC 開始時に「3 つのインプット」をご提供いただき、', 'size': 11, 'color': NAVY, 'bold': True},
        {'text': '　 初回 3C レポートを 1 営業日で出力。', 'size': 11, 'color': NAVY, 'bold': True},
    ])

    # 下部脚注：スクレイピング 101 件は実 EC データ（NAVY バー）
    add_lines_box(s, Inches(0.4), Inches(5.75), Inches(12.5), Inches(1.15),
                  fill=NAVY, border=NAVY, border_w=0.75, lines=[
        {'text': '【スクレイピングデータ 101 件は実 EC データ】', 'size': 12, 'bold': True, 'color': WHITE},
        {'text': '内訳：asone 25 / トミー精工 12 / ヤマト科学 14 / 平山 27 / アルプ 23（2026-05 取得）', 'size': 10, 'color': WHITE},
        {'text': '本提案書 Slide 8 の競合スペック比較表・Slide 10 の価格ポジショニングはこの実データに基づく。', 'size': 10, 'color': WHITE},
        {'text': '→ サンプル（POS/SNS）と実データ（スクレイピング 101 件）の線引きを明確化', 'size': 10, 'bold': True, 'color': ACCENT_BG},
    ])


def main():
    prs = Presentation(str(SRC))
    counts = {old: 0 for old, _ in REPLACEMENTS}
    slides = list(prs.slides)
    print(f"loaded V6: {len(slides)} slides")

    # Find/Replace
    for slide in slides:
        for tf in _walk_text_frames(slide):
            for para in tf.paragraphs:
                for old, new in REPLACEMENTS:
                    if _replace_in_runs(para, old, new):
                        counts[old] += 1

    # 新規スライド追加
    slide_21_data_premise(prs)
    print(f"after add: {len(prs.slides)} slides")

    prs.save(str(DST))
    print(f"saved V7: {DST}")
    print("\n--- replacement counts ---")
    for old, c in counts.items():
        marker = '  ' if c > 0 else '  [MISS] '
        print(f"{marker}[{c}] {old[:60]}{'...' if len(old) > 60 else ''}")


if __name__ == "__main__":
    main()
