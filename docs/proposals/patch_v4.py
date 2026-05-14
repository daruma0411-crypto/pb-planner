"""V3.pptx に 2 枚追加 → V4.pptx として保存

追加:
- Slide A: ハルシネーション対策（事実担保の多層ガード）
- Slide B: 統括エージェントとの壁打ち（チャットモード ユースケース）
"""
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V3.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V4.pptx"


# build_pptx.py と同じ色・フォント
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_LIGHT = RGBColor(0x3A, 0x5A, 0x85)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)
ACCENT = RGBColor(0xC4, 0x76, 0x3A)
ACCENT_BG = RGBColor(0xFD, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x18, 0x27)
FONT = 'Noto Sans JP'

TITLE_PREFIX = '勘と手作業から脱却する、次世代PB開発の「設計図」'
FOOTER_TXT = 'アズワン株式会社様向け｜オートクレーブPB戦略提案｜2026年5月'


def _set_font(run, size=14, bold=False, color=BLACK, name=FONT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, border=None,
             border_w=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(border_w)
    else:
        tb.line.fill.background()
    return tb


def add_lines_box(slide, x, y, w, h, lines, fill=None, border=None, border_w=1.0):
    """複数行のテキストボックス。lines は [{text, size, bold, color}, ...]"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln.get('text', '')
        _set_font(run,
                  size=ln.get('size', 12),
                  bold=ln.get('bold', False),
                  color=ln.get('color', BLACK))
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(border_w)
    else:
        tb.line.fill.background()
    return tb


def add_title(slide, title):
    add_text(slide, Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.4),
             TITLE_PREFIX, size=11, color=GRAY)
    add_text(slide, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.55),
             title, size=22, bold=True, color=NAVY)


def add_page_meta(slide, page_no, total=20):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8.0), Inches(0.3),
             FOOTER_TXT, size=9, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
             f'{page_no} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def slide_19_hallucination(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_meta(s, 19)
    add_title(s, '事実に立脚するレポートを担保する「多層ガード」')

    # 左：主軸の DB
    add_text(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.4),
             '【主軸】一次情報のスクレイピング DB', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(0.4), Inches(1.95), Inches(6.2), Inches(2.3),
                  fill=GRAY_BG, border=GRAY, border_w=0.75, lines=[
        {'text': '・自社（AS ONE / ナビス AXEL）：25 機種', 'size': 12, 'bold': True},
        {'text': '・製造パートナー（トミー精工）：12 機種', 'size': 12},
        {'text': '・競合（ヤマト科学／平山製作所／アルプ）：64 機種', 'size': 12},
        {'text': '─ 合計 101 件、価格・型番・スペック・URL を実在ベースで保持 ─', 'size': 12, 'bold': True, 'color': ACCENT},
        {'text': '', 'size': 4},
        {'text': '→ TOP 機種スペック比較表・価格対決・共食い検証は', 'size': 11},
        {'text': '  この実 DB から引用、ハルシネーション余地ゼロ', 'size': 11, 'color': NAVY, 'bold': True},
    ])

    # 右：プロンプトレベルの対策
    add_text(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.4),
             '【プロンプト】Claude への 4 つの厳格指示', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(6.8), Inches(1.95), Inches(6.2), Inches(2.3),
                  fill=ACCENT_BG, border=ACCENT, border_w=0.75, lines=[
        {'text': '① 「ハルシネーション厳禁」明示', 'size': 12, 'bold': True},
        {'text': '② 「データに記載なし」を許容＋明記必須', 'size': 12, 'bold': True},
        {'text': '③ 各事実に出典セクション付与', 'size': 12, 'bold': True},
        {'text': '　（3C／POS／SNS／競合 のいずれか）', 'size': 11, 'color': GRAY},
        {'text': '④ 訓練データ由来は「推計／推定」と自己ラベリング', 'size': 12, 'bold': True},
        {'text': '', 'size': 4},
        {'text': '→ Claude が honest に「データに記載なし」と書く', 'size': 11, 'color': NAVY, 'bold': True},
    ])

    # 下：多層構造の図解（横長）
    add_text(s, Inches(0.4), Inches(4.45), Inches(12.5), Inches(0.4),
             '【パイプライン構造】各 phase が前段レポートを必須引用', size=14, bold=True, color=NAVY)

    phases = [
        ('① 3C\n（事実忠実）', 'スクレイピング DB を一次引用'),
        ('② KSF\n（出典必須）', '3C＋POS＋SNS から帰納'),
        ('③ STP\n（KSF 駆動）', 'KSF を根拠にセグメント定義'),
        ('④ 4P\n（STP 駆動）', 'STP のターゲットを実行戦術化'),
        ('⑤ キャッチ\n＋マスタ', '全段の論理連鎖を集約'),
    ]
    bw = Inches(2.4)
    bx = Inches(0.4)
    for i, (label, desc) in enumerate(phases):
        x = bx + bw * i + Inches(0.05 * i)
        add_lines_box(s, x, Inches(4.95), bw, Inches(1.3),
                      fill=NAVY, border=NAVY, border_w=1.0, lines=[
            {'text': label, 'size': 12, 'bold': True, 'color': WHITE},
            {'text': '', 'size': 4},
            {'text': desc, 'size': 9, 'color': WHITE},
        ])

    # 下部脚注：今後の強化ロードマップ
    add_lines_box(s, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.65),
                  fill=GRAY_BG, border=GRAY, border_w=0.5, lines=[
        {'text': '【さらに強化するロードマップ】', 'size': 10, 'bold': True, 'color': NAVY},
        {'text': '行レベル出典トラッキング（クリック→ raw data へジャンプ） ／ Web 検索 API 接続でファクト強化 ／ ファクトチェック専用エージェント追加', 'size': 9, 'color': BLACK},
    ])


def slide_20_chat(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_meta(s, 20)
    add_title(s, 'レポート完成後の壁打ち：統括エージェントとの対話モード')

    # 左：仕組み
    add_text(s, Inches(0.4), Inches(1.5), Inches(5.5), Inches(0.4),
             '【仕組み】統括エージェントが全コンテキストを把握', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(0.4), Inches(1.95), Inches(5.5), Inches(4.7),
                  fill=GRAY_BG, border=GRAY, border_w=0.75, lines=[
        {'text': '案件詳細画面に常設のチャット UI を配置。', 'size': 12},
        {'text': '統括コンサルタント AI が以下をすべて把握:', 'size': 12},
        {'text': '', 'size': 4},
        {'text': '・案件メタ（カテゴリ・PB コンセプト・ベース機種）', 'size': 11},
        {'text': '・スクレイピング 101 件の生データ', 'size': 11},
        {'text': '・POS / SNS の手入力サマリ', 'size': 11},
        {'text': '・5 フェーズ全レポート Markdown 全文', 'size': 11, 'bold': True, 'color': ACCENT},
        {'text': '', 'size': 4},
        {'text': '→ 「あの数値の出典は？」「あの軸で再検討して」', 'size': 11, 'color': NAVY, 'bold': True},
        {'text': '　 のような対話的ブラッシュアップが可能', 'size': 11, 'color': NAVY, 'bold': True},
        {'text': '', 'size': 4},
        {'text': '【履歴】案件単位で永続化、いつでも会話を再開', 'size': 11, 'bold': True, 'color': NAVY},
        {'text': '', 'size': 4},
        {'text': '【ルール】レポート未記載は「データに記載なし」と返答、', 'size': 11},
        {'text': '　 ハルシネーション抑止は生成時と同等', 'size': 11},
    ])

    # 右：ユースケース
    add_text(s, Inches(6.1), Inches(1.5), Inches(6.9), Inches(0.4),
             '【ユースケース】よくある質問・指示パターン', size=14, bold=True, color=NAVY)

    cases = [
        ('🔍 ファクト確認', '「Customer 1-1 の市場規模 12,000 台の出典は？」', 'スクレイピング DB か手入力サマリか訓練データ由来かを明示'),
        ('📊 根拠掘り下げ', '「KSF #3 のペイン強度『高』の判定根拠は？」', '該当データ出典＋類似事例の参照'),
        ('♻️ 軸変更・再生成', '「STP の T1 ターゲット、価格軸 → 購買サイクル軸で再検討して」', '該当 phase の再生成を提案＋実行（next phase）'),
        ('✂️ コピー調整', '「キャッチ C2 を 30 文字以内に短縮、女性向け強調」', '4P の訴求軸を維持しつつ短縮案を 3 つ提示'),
        ('📄 アウトプット', '「上長会議用に 1 ページサマリを Word で出力」', 'PIM/PPTX/Word 自動生成（既存ツール連動）'),
    ]
    y = Inches(1.95)
    for icon_title, q, a in cases:
        add_lines_box(s, Inches(6.1), y, Inches(6.9), Inches(0.92),
                      fill=ACCENT_BG, border=ACCENT, border_w=0.5, lines=[
            {'text': icon_title, 'size': 11, 'bold': True, 'color': ACCENT},
            {'text': f'Q：{q}', 'size': 10, 'color': BLACK},
            {'text': f'A：{a}', 'size': 9.5, 'color': GRAY},
        ])
        y += Inches(0.97)

    # 下部脚注
    add_lines_box(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.4),
                  fill=NAVY, border=NAVY, border_w=0.5, lines=[
        {'text': '一発生成 × 対話的ブラッシュアップ ＝ 担当者の「考える時間」を最大化', 'size': 12, 'bold': True, 'color': WHITE},
    ])


def main():
    prs = Presentation(str(SRC))
    print(f"loaded: {SRC}")
    print(f"  existing slides: {len(prs.slides)}")
    slide_19_hallucination(prs)
    slide_20_chat(prs)
    print(f"  after add: {len(prs.slides)}")
    prs.save(str(DST))
    print(f"saved: {DST}")


if __name__ == "__main__":
    main()
