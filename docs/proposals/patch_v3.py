"""V2.pptx を実装内容に合わせて修正 → V3.pptx として保存"""
from pathlib import Path
from pptx import Presentation


SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V2.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V3.pptx"


# 実装と提案書の乖離を埋める Find/Replace マップ
REPLACEMENTS = [
    # Slide 3: リードタイム表現を実証ベースに
    ("圧倒的なリードタイム短縮（数日〜数週）",
     "圧倒的なリードタイム短縮（実証：7-10 分/案件）"),

    # Slide 4: 5 フェーズ表記に修正（実装は KSF + キャッチ+マスタも自動化済）
    ("3C → KSF→STP → 4P を自動生成",
     "3C → KSF → STP → 4P → キャッチ+マスタ を自動生成"),
    ("3C → STP → 4P を自動生成",
     "3C → KSF → STP → 4P → キャッチ+マスタ を自動生成"),

    # Slide 11: アウトプットを実装通りに（PDF API は撤退、HTML/MD/ブラウザ印刷）
    ("PB企画書PDF", "PB企画書 HTML"),
    ("3C／KSF／STP／4P／マスタを統合",
     "3C／KSF／STP／4P／キャッチ+マスタを統合（ブラウザ印刷で PDF 化）"),

    # Slide 12: 3C もペルソナ化済み
    ("（既存3Cレポート機能）", "シニア戦略コンサル"),

    # Slide 13: 圧縮倍率を実証値に
    ("※ アズワン社の既存PB企画フロー（典型6ヶ月）を、AI支援により 1/4〜1/8 に圧縮できる想定。",
     "※ アズワン社の既存PB企画フロー（典型6ヶ月）を、AI支援により 1/100〜1/1000 程度に圧縮（実証：7-10 分/案件、MVP 稼働中）。"),

    # Slide 15: 5 営業日 → 1 日（MVP 稼働中）
    ("これら 3 つのインプットがそろい次第、初回アウトプット（3Cレポート）まで【最短 5 営業日】で提示可能。",
     "これら 3 つのインプットがそろい次第、初回アウトプット（3Cレポート）まで【最短 1 営業日】で提示可能（MVP 稼働中）。"),

    # Slide 7: 「3C → STP → 4P」を 5 フェーズに修正
    ("膨大な情報を数分で「3C → STP → 4P」のマーケティング・フレームワークへ【自動構造化】。",
     "膨大な情報を数分で「3C → KSF → STP → 4P → キャッチ+マスタ」の 5 段マーケティング・フレームワークへ【自動構造化】。"),

    # Slide 17: 「マスタ生成まで」を明示
    ("情報収集から 3C → KSF → STP → 4P、マスタ生成までを",
     "情報収集から 3C → KSF → STP → 4P → キャッチ+マスタ までを"),

    # Slide 3 表: 「4専門家AIエージェント」→ 5 専門家
    ("4専門家AIエージェントの分業＋ヒト承認",
     "5専門家AIエージェント（3C／KSF／STP／4P／キャッチ+マスタ）の分業＋ヒト承認"),

    # Next Action: タイミングを実証ベースで更新
    ("〜5週間", "〜1週間（MVP 稼働中、既存環境に接続のみ）"),
]


def _replace_in_runs(paragraph, old, new):
    """run.text レベルで置換。run をまたぐ場合は paragraph 全体を 1 run に潰す。"""
    full_text = "".join(run.text for run in paragraph.runs)
    if old not in full_text:
        return False
    new_full = full_text.replace(old, new)
    if not paragraph.runs:
        return False
    # 最初の run に全文置き、他の run は空に
    first = paragraph.runs[0]
    first.text = new_full
    for run in paragraph.runs[1:]:
        run.text = ""
    return True


def _walk_text_frames(slide):
    """slide 内の全 text_frame を yield"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        # group の中も降りる
        if hasattr(shape, 'shapes'):
            for s in shape.shapes:
                if hasattr(s, 'text_frame') and s.has_text_frame:
                    yield s.text_frame


def main():
    prs = Presentation(str(SRC))
    counts = {old: 0 for old, _ in REPLACEMENTS}
    for slide in prs.slides:
        for tf in _walk_text_frames(slide):
            for para in tf.paragraphs:
                for old, new in REPLACEMENTS:
                    if _replace_in_runs(para, old, new):
                        counts[old] += 1
    prs.save(str(DST))
    print(f"saved: {DST}")
    print("--- replacement counts ---")
    for old, c in counts.items():
        print(f"  [{c}] {old[:50]}{'...' if len(old) > 50 else ''}")


if __name__ == "__main__":
    main()
