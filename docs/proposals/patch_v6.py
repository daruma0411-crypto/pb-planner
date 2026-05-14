"""V4 → V6: 赤字注記なし版、本文を「FLS-1000 弱点補強」軸で書き換え

新軸 (新 STP/4P/Finish レポートから抽出):
- 主軸 T-A: 製薬・CDMO 試作ラボ向け = IQ-OQ + USB/LAN データ可搬 + 電子署名
- 副軸 T-C: バイオベンチャー/再生医療向け = 多言語UI + クラウド通知 + サブスク
- woman-friendly は NB（FLS-1000）の訴求軸として継承、PB 独自軸からは外す
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V4.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V6.pptx"


REPLACEMENTS = [
    # Slide 6: PB 空白カテゴリ仮説
    ("100L クラス × 人間工学",
     "100L × GMP/IQ-OQ + データ可搬"),
    ("（女性研究者・若手職員配慮）",
     "（製薬・CDMO 試作ラボ向け）"),
    ("100L クラス × 設置容易性",
     "100L × 多言語UI + クラウド通知 + サブスク"),
    ("（単相100V / スリム筐体）",
     "（バイオベンチャー・再生医療向け）"),
    ("100L クラス × 価格訴求",
     "100L × 横置き省設置 + 消耗品同梱"),
    ("（NB比 -25〜30%、機能限定版）",
     "（量販向け、IQ-OQ 込み総コストで -25%）"),

    # Slide 7: ポジショニング・マトリクス
    ("人間工学×中価格",
     "IQ-OQ×中価格"),
    ("woman-friendly × 100L",
     "validated 100L (IQ-OQ + データ可搬)"),
    ("ボリュームゾーン",
     "NB ボリュームゾーン"),

    # Slide 9: スペックちょい替え PB 仮称
    ("woman-friendly 100L（PB 仮称）",
     "validated 100L（PB 仮称：IQ-OQ + データ可搬）"),
    ("・軽量化蓋（電動アシスト）",
     "・IQ-OQ 文書一式同梱（QA 負担ゼロ化）"),
    ("・低位置投入口（150mm 短縮）",
     "・USB/LAN データ出力 + 電子署名対応"),
    ("・カラータッチパネル（多言語）",
     "・多言語UI（日英中韓）+ クラウド通知"),
    ("・「重い蓋から、解放を。」",
     "・「IQ-OQ、最初から動く 100L」"),
    ("・「100L、ラボの誰もが扱えるサイズへ」",
     "・「QA の負担を、買った瞬間に終わらせる」"),
    ("・NB 比 -10〜15%、機能維持",
     "・NB 比 -5〜10%、IQ-OQ 込みで総コスト -25%"),

    # Slide 10: 想定スペック・キャッチ・価格・ターゲット
    ('PB「woman-friendly 100L」想定スペック・訴求・価格',
     'PB「validated 100L」想定スペック・訴求・価格'),
    ("電動アシスト上蓋",
     "IQ-OQ 文書 同梱"),
    ("床から850mm",
     "USB/LAN/クラウド対応"),
    ("7インチカラータッチ",
     "21CFR Part11 電子署名対応"),
    ("標準（短時間モード）",
     "標準（短時間モード）+ クラウド通知"),
    ('① 「重い蓋から、解放を。」',
     '① 「IQ-OQ、最初から動く 100L」'),
    ('　 ── 電動アシストで誰でも安全に',
     '　 ── 文書一式同梱、QA 負担をゼロ化'),
    ('② 「100L、ラボの誰もが扱えるサイズへ。」',
     '② 「QA の負担を、買った瞬間に終わらせる」'),
    ('③ 「やさしいオートクレーブ、はじまる。」',
     '③ 「データは、研究室を出てクラウドへ」'),
    ("・想定価格：NB（FLS-1000）比 -10〜15%",
     "・想定価格：NB 比 -5〜10%（IQ-OQ込み総コストで -25%）"),
    ("・上位機（平山HV）比 -25〜35%",
     "・GMP対応上位機（平山HV）比 -25〜35%"),
    ("・機能維持＋人間工学プレミアム",
     "・NB 同等の人間工学＋IQ-OQ・データ可搬で総コスト優位"),
    ("若手女性研究者（D2〜ポスドク）／ラボマネージャー（共用機運用）",
     "T-A：製薬・CDMO 試作ラボ QA 担当（40代）／T-C：バイオベンチャー研究 PI（再生医療・治験原料）"),
    ("AXEL 主力カタログ掲載／大学法人ルート（共同購入）／研究機器商社との協業",
     "AXEL 主力カタログ＋アズワン直販ルート／製薬向け代理店経由（IQ-OQ パッケージ訴求）／消耗品サブスク同梱"),

    # その他「woman-friendly」表記がもしあれば
    ("woman-friendly", "validated 100L"),
]


# Speaker Notes に追加する補足
NOTES = {
    5: """【V6 修正の背景】
当初の仮説「100L × 人間工学（woman-friendly）」は、ベース機種候補のトミー精工 FLS-1000 が既に NB として「女性ユーザー着目・片手開閉アシスト・ローテーブル79cm設計」を訴求済みであることが判明。
そのため、PB の独自差別化軸を「FLS-1000 の弱点補強」軸に再設計した。
- 仮説 A: GMP/IQ-OQ + USB/LAN データ可搬（製薬・CDMO 試作ラボ向け）
- 仮説 B: 多言語UI + クラウド通知 + サブスク（バイオベンチャー/再生医療向け）
- 仮説 C: 横置き省設置 + 消耗品同梱（量販向け）
woman-friendly 軸は NB が占有、PB は継承するが独自軸にしない。""",
    6: """【V6 修正】ポジショニング・マトリクスの空白地帯を「IQ-OQ × 中価格」に再定義。
新 STP レポートの軸：
- 軸X：滅菌記録の外部要件強度（IQ-OQ／GMP）
- 軸Y：ハード単価感度 × 消耗品月次回転
優先ターゲット T-A（製薬 CDMO 試作）、T-C（バイオベンチャー）は NB 全社が現状 satisfy していない空白地帯。""",
    8: """【V6 修正】PB 仮称を「validated 100L」に変更。
NB（FLS-1000）の人間工学訴求は維持しつつ、PB の独自軸として以下を加算：
1. IQ-OQ 文書一式同梱（業界全社が現状未訴求）
2. USB/LAN データ可搬 + 電子署名（21CFR Part11 対応）
3. 多言語 UI + クラウド通知
4. NB 比 -5〜10% 価格、IQ-OQ 込み総コストで -25%
キャッチコピーも NB 軸（woman-friendly）の言い換えから、IQ-OQ・データ可搬・クラウド連動の文脈に刷新。""",
    9: """【V6 修正】想定スペック・キャッチ・価格・ターゲット・販路を新軸で再設計。
ターゲット T-A：製薬・CDMO 試作ラボの QA 担当（40代、IQ-OQ × データ可搬を重視）
ターゲット T-C：バイオベンチャー研究 PI（多言語 + クラウド + サブスクで未開拓層を取りに行く）
T-B（大学）は NB 競合度が高い（hirayama HV-110Ⅱ ¥990,000、yamato HVA-110LB ¥742,000）ため優先度を下げる。
詳細は新 STP/4P/Finish レポート参照（HTML 別添）。""",
}


def _replace_in_runs(paragraph, old, new):
    full_text = "".join(run.text for run in paragraph.runs)
    if old not in full_text:
        return False
    new_full = full_text.replace(old, new)
    if not paragraph.runs:
        return False
    first = paragraph.runs[0]
    first.text = new_full
    for run in paragraph.runs[1:]:
        run.text = ""
    return True


def _walk_text_frames(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape.text_frame
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        if hasattr(shape, 'shapes'):
            for s in shape.shapes:
                if hasattr(s, 'text_frame') and s.has_text_frame:
                    yield s.text_frame


def set_notes(slide, text):
    """Speaker Notes を設定"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = text


def main():
    prs = Presentation(str(SRC))
    counts = {old: 0 for old, _ in REPLACEMENTS}
    slides = list(prs.slides)
    print(f"loaded V4: {len(slides)} slides")

    for slide in slides:
        for tf in _walk_text_frames(slide):
            for para in tf.paragraphs:
                for old, new in REPLACEMENTS:
                    if _replace_in_runs(para, old, new):
                        counts[old] += 1

    # Speaker Notes 追加
    for idx, note_text in NOTES.items():
        if idx < len(slides):
            set_notes(slides[idx], note_text)
            print(f"  speaker notes set on slide {idx + 1}")

    prs.save(str(DST))
    print(f"saved V6: {DST}")
    print("\n--- replacement counts ---")
    for old, c in counts.items():
        print(f"  [{c}] {old[:60]}{'...' if len(old) > 60 else ''}")


if __name__ == "__main__":
    main()
