"""V7 → V8: 費用感スライドを最終ページに追加"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


SRC = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V7.pptx"
DST = Path(__file__).parent / "Precision_PB_Blueprint_Aswan_editable_V8.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xC4, 0x76, 0x3A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)
ACCENT_BG = RGBColor(0xFD, 0xF6, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x18, 0x27)
FONT = 'Noto Sans JP'


def _set_font(run, size=14, bold=False, color=BLACK):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, border=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(1.0)
    else:
        tb.line.fill.background()
    return tb


def add_lines_box(slide, x, y, w, h, lines, fill=None, border=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln.get('text', '')
        _set_font(run,
                  size=ln.get('size', 12),
                  bold=ln.get('bold', False),
                  color=ln.get('color', BLACK))
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(1.0)
    else:
        tb.line.fill.background()
    return tb


def slide_cost(prs, page_no, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # ヘッダ
    add_text(s, Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.4),
             '勘と手作業から脱却する、次世代PB開発の「設計図」', size=11, color=GRAY)
    add_text(s, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.55),
             '費用感とフェーズ別投資イメージ', size=22, bold=True, color=NAVY)

    # 左：PoC 期間
    add_text(s, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.4),
             '【PoC 期間 ─ 初期】', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(0.4), Inches(1.95), Inches(6.2), Inches(3.0),
                  fill=GRAY_BG, border=NAVY, lines=[
        {'text': '・現プロダクト利用 ……………… ¥0', 'size': 13, 'bold': True},
        {'text': '   ベース機能（5 フェーズパイプライン）は既に稼働中',
         'size': 10, 'color': GRAY},
        {'text': '', 'size': 4},
        {'text': '・要件カスタマイズ ………… 約 ¥200,000（一括）',
         'size': 13, 'bold': True, 'color': ACCENT},
        {'text': '   貴社固有データ連携・UI 調整・帳票テンプレ・追加スクレイパー 等',
         'size': 10, 'color': GRAY},
        {'text': '', 'size': 4},
        {'text': '・想定期間：1〜2 ヶ月', 'size': 12, 'color': BLACK},
    ])

    # 右：ランニング
    add_text(s, Inches(6.8), Inches(1.5), Inches(6.2), Inches(0.4),
             '【以降ランニング ─ 月次】', size=14, bold=True, color=NAVY)
    add_lines_box(s, Inches(6.8), Inches(1.95), Inches(6.2), Inches(3.0),
                  fill=ACCENT_BG, border=ACCENT, lines=[
        {'text': '・データ更新（競合 + 対象製品 100 件）  約 ¥30,000',
         'size': 13, 'bold': True},
        {'text': '   月次スクレイピング再取得・差分反映',
         'size': 10, 'color': GRAY},
        {'text': '', 'size': 4},
        {'text': '・保守サービス ………………… ¥50,000 / 月',
         'size': 13, 'bold': True},
        {'text': '   リモート障害対応・モデル更新・プロンプト改善',
         'size': 10, 'color': GRAY},
        {'text': '', 'size': 4},
        {'text': '・LLM API 実費（想定）……… 約 ¥10,000 / 月',
         'size': 13, 'bold': True},
        {'text': '   利用ボリュームに応じて変動（実費精算）',
         'size': 10, 'color': GRAY},
        {'text': '', 'size': 4},
        {'text': '─────────────────────────', 'size': 10, 'color': GRAY},
        {'text': '【月次合計】 約 ¥90,000', 'size': 14, 'bold': True, 'color': ACCENT},
    ])

    # 下：年間総額試算
    add_text(s, Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.35),
             '【年間総額試算（参考）】', size=13, bold=True, color=NAVY)
    add_lines_box(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.0),
                  fill=NAVY, border=NAVY, lines=[
        {'text': '初年度：PoC ¥200,000 ＋ 月次 ¥90,000 × 12 = 約 ¥1,280,000',
         'size': 13, 'bold': True, 'color': WHITE},
        {'text': '2 年目以降：月次 ¥90,000 × 12 = 約 ¥1,080,000 / 年',
         'size': 13, 'bold': True, 'color': WHITE},
    ])

    # 注記
    add_lines_box(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.4),
                  fill=GRAY_BG, lines=[
        {'text': '※ 全て税別／LLM API は実費精算（変動）／保守はリモート対応／追加要件は別途見積',
         'size': 10, 'color': GRAY},
    ])

    # ページメタ
    add_text(s, Inches(0.4), Inches(7.05), Inches(8.0), Inches(0.3),
             'アズワン株式会社様向け｜オートクレーブPB戦略提案｜2026年5月',
             size=9, color=GRAY)
    add_text(s, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
             f'{page_no} / {total}', size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def main():
    prs = Presentation(str(SRC))
    n = len(prs.slides)
    print(f"loaded V7: {n} slides")
    slide_cost(prs, page_no=n + 1, total=n + 1)
    print(f"after add: {len(prs.slides)} slides")
    prs.save(str(DST))
    print(f"saved V8: {DST}")


if __name__ == "__main__":
    main()
