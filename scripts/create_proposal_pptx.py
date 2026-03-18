"""
PB企画プランナー 提案書パワーポイント生成スクリプト
構成:
  1. コンセプトページ
  2. フローページ
  3-7. フレームワーク分析（3C, SWOT, ポジショニング, 5Forces, 価格帯マップ）
  8. PBカード完成
  9. 仕様差分作成
  10. 企画書Word生成
  11. マスタExcel完成イメージ
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SS_DIR = r"C:\Users\iwashita.AKGNET\Pictures\Screenshots"
OUTPUT = r"C:\Users\iwashita.AKGNET\pb-planner\docs\PB企画プランナー_提案書.pptx"

# --- Color palette ---
RED_MAIN = RGBColor(0xDC, 0x26, 0x26)       # メインレッド
RED_DARK = RGBColor(0xB9, 0x1C, 0x1C)       # ダークレッド
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BLUE = RGBColor(0x21, 0x96, 0xF3)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)

# Slide dimensions: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_title_bar(slide, title_text, subtitle_text=None):
    """赤帯タイトルバーを上部に配置"""
    # 赤帯背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_W, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED_MAIN
    shape.line.fill.background()

    # タイトルテキスト
    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15), Inches(12), Inches(0.6)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans JP"

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.7), Inches(12), Inches(0.4)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
        p2.font.name = "Noto Sans JP"


def add_footer(slide, page_num, total=11):
    """フッター"""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(12), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"PB企画プランナー 提案書  |  AKGNET  |  {page_num}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.RIGHT


def add_screenshot_slide(prs, title, subtitle, img_filename, page_num, caption=None):
    """スクリーンショットスライドテンプレート"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    add_title_bar(slide, title, subtitle)

    img_path = os.path.join(SS_DIR, img_filename)
    if os.path.exists(img_path):
        # 画像を中央に大きく配置（左右マージン0.4in、上は赤帯の下）
        slide.shapes.add_picture(
            img_path,
            Inches(0.4), Inches(1.4),
            Inches(12.5), Inches(5.4)
        )
    else:
        # 画像なしプレースホルダー
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = GRAY
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[{img_filename}]"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    if caption:
        txBox = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.85), Inches(12), Inches(0.3)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.font.name = "Noto Sans JP"

    add_footer(slide, page_num)
    return slide


def add_box(slide, left, top, width, height, fill_color, text, font_size=14, text_color=BLACK, bold=False):
    """テキストボックスを追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)
    return shape


def add_arrow(slide, left, top, width=Inches(0.6), height=Inches(0.01)):
    """矢印"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        left, top, width, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED_MAIN
    shape.line.fill.background()
    return shape


def create_slide_cover(prs):
    """表紙"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景を赤に
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RED_MAIN

    # メインタイトル
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(11), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PB企画プランナー"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.CENTER

    # サブタイトル
    txBox2 = slide.shapes.add_textbox(
        Inches(1), Inches(3.0), Inches(11), Inches(1.0)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "AIを活用したPB企画支援プラットフォーム"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0xFF, 0xDD, 0xDD)
    p2.font.name = "Noto Sans JP"
    p2.alignment = PP_ALIGN.CENTER

    # 区切り線
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4), Inches(4.2), Inches(5), Inches(0.03)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()

    # 会社名
    txBox3 = slide.shapes.add_textbox(
        Inches(1), Inches(4.6), Inches(11), Inches(0.6)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "株式会社あかがね"
    p3.font.size = Pt(20)
    p3.font.color.rgb = WHITE
    p3.font.name = "Noto Sans JP"
    p3.alignment = PP_ALIGN.CENTER

    # 日付
    txBox4 = slide.shapes.add_textbox(
        Inches(1), Inches(5.3), Inches(11), Inches(0.5)
    )
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "2026年3月"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0xFF, 0xBB, 0xBB)
    p4.font.name = "Noto Sans JP"
    p4.alignment = PP_ALIGN.CENTER


def create_slide_concept(prs):
    """P1: コンセプトページ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide, "コンセプト", "企画書を作成しながら商品マスタ作成が同時に可能")

    # 3つの柱
    col_w = Inches(3.5)
    col_h = Inches(4.0)
    gap = Inches(0.5)
    start_x = Inches(0.8)
    top_y = Inches(1.8)

    items = [
        {
            "icon": "AI",
            "title": "AI企画書作成",
            "color": BLUE,
            "desc": "生成AIによる\n企画書作成工数削減",
            "detail": "対話形式でPB企画を進行\nフレームワーク分析を自動実行\n企画書(Word)を自動生成"
        },
        {
            "icon": "DB",
            "title": "競合調査代行",
            "color": GREEN,
            "desc": "インプットデータ：\n競合商品情報収集の代行",
            "detail": "仕入先・競合製品データを\nDB登録（当社が代行）\n76製品の実データで分析"
        },
        {
            "icon": "PIM",
            "title": "後工程同時完了",
            "color": ORANGE,
            "desc": "後工程部門に任せていた\n作成業務を同時完了",
            "detail": "PIMデータ(Excel)自動生成\nJANコード・品番管理\n商品マスタ登録準備完了"
        }
    ]

    for i, item in enumerate(items):
        x = start_x + (col_w + gap) * i

        # アイコン円
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + Inches(1.1), top_y, Inches(1.2), Inches(1.2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = item["color"]
        circle.line.fill.background()
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = item["icon"]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Noto Sans JP"
        p.alignment = PP_ALIGN.CENTER

        # タイトル
        txBox = slide.shapes.add_textbox(
            x, top_y + Inches(1.4), col_w, Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = item["title"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.font.name = "Noto Sans JP"
        p.alignment = PP_ALIGN.CENTER

        # 説明（赤テキスト）
        txBox2 = slide.shapes.add_textbox(
            x, top_y + Inches(1.95), col_w, Inches(0.9)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = item["desc"]
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = RED_MAIN
        p2.font.name = "Noto Sans JP"
        p2.alignment = PP_ALIGN.CENTER

        # 詳細
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.1), top_y + Inches(2.9), col_w - Inches(0.2), Inches(1.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()
        tf3 = box.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = item["detail"]
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRAY
        p3.font.name = "Noto Sans JP"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(6)

    # 底部キャッチ
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(6.3), Inches(11), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "企画立案から商品マスタ作成まで、ワンストップで完了"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_MAIN
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 1)


def create_slide_flow(prs):
    """P2: フローページ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_bar(slide, "ご利用フロー", "DB登録から企画書完成まで最短2日")

    # フロー3ステップ
    step_w = Inches(3.2)
    step_h = Inches(3.8)
    start_x = Inches(0.5)
    top_y = Inches(1.6)
    arrow_w = Inches(0.6)

    steps = [
        {
            "num": "STEP 1",
            "title": "プラットフォーム提供",
            "who": "あかがね",
            "color": RED_MAIN,
            "items": [
                "PB企画プランナー環境構築",
                "LLM（Claude AI）連携",
                "製品データベース構築",
                "チャット対話型UI",
                "ドキュメント自動生成",
            ]
        },
        {
            "num": "STEP 2",
            "title": "競合・対象製品DB登録",
            "who": "あかがね（1日）",
            "color": BLUE,
            "items": [
                "対象製品のスペック収集",
                "競合メーカー製品スクレイピング",
                "JSON形式でDB登録",
                "価格・スペック・特徴整理",
                "例: オートクレーブ4社76製品",
            ]
        },
        {
            "num": "STEP 3",
            "title": "PB企画プランナー利用",
            "who": "貴社（最短1日）",
            "color": GREEN,
            "items": [
                "チャットでPB企画を対話的に進行",
                "5種のフレームワーク分析",
                "PB仕様を決定・差分記録",
                "企画書(Word)自動生成",
                "PIMデータ(Excel)同時完了",
            ]
        }
    ]

    for i, step in enumerate(steps):
        x = start_x + (step_w + arrow_w + Inches(0.15)) * i

        # STEP番号バッジ
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.8), top_y, Inches(1.6), Inches(0.45)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = step["color"]
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = step["num"]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Noto Sans JP"
        p.alignment = PP_ALIGN.CENTER

        # タイトル
        txBox = slide.shapes.add_textbox(
            x, top_y + Inches(0.6), step_w, Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = step["title"]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.font.name = "Noto Sans JP"
        p.alignment = PP_ALIGN.CENTER

        # 担当
        txBox2 = slide.shapes.add_textbox(
            x, top_y + Inches(1.05), step_w, Inches(0.35)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"担当: {step['who']}"
        p2.font.size = Pt(11)
        p2.font.color.rgb = step["color"]
        p2.font.bold = True
        p2.font.name = "Noto Sans JP"
        p2.alignment = PP_ALIGN.CENTER

        # 内容リスト
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.15), top_y + Inches(1.5), step_w - Inches(0.3), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.fill.background()

        tf3 = box.text_frame
        tf3.word_wrap = True
        for j, item in enumerate(step["items"]):
            if j == 0:
                p = tf3.paragraphs[0]
            else:
                p = tf3.add_paragraph()
            p.text = f"  {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = BLACK
            p.font.name = "Noto Sans JP"
            p.space_before = Pt(4)

        # 矢印（最後以外）
        if i < 2:
            arrow_x = x + step_w + Inches(0.05)
            add_arrow(slide, arrow_x, top_y + Inches(2.0))

    # 下部メッセージ
    msg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(5.8), Inches(10), Inches(0.8)
    )
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    msg_box.line.color.rgb = RED_MAIN
    msg_box.line.width = Pt(1.5)

    tf = msg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DB登録（1日）+ 企画作成（最短1日）= 最短2日で企画書 + 商品マスタ完成"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED_MAIN
    p.font.name = "Noto Sans JP"
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 2)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 表紙
    create_slide_cover(prs)

    # P1: コンセプト
    create_slide_concept(prs)

    # P2: フロー
    create_slide_flow(prs)

    # P3-7: フレームワーク分析スクリーンショット
    ss_slides = [
        {
            "title": "3C分析",
            "subtitle": "Company / Competitor / Customer の3視点で市場を分析",
            "img": "pb-planner-03-3c.png",
            "caption": "AIがDB上の76製品データを参照し、3C分析を自動生成"
        },
        {
            "title": "SWOT分析",
            "subtitle": "強み・弱み・機会・脅威を可視化",
            "img": "pb-planner-04-swot-bottom.png",
            "caption": "2x2カラーグリッドでSWOT分析を直感的に表示"
        },
        {
            "title": "ポジショニング分析",
            "subtitle": "競合マップで空白ゾーンを発見",
            "img": "pb-planner-08-positioning-chart.png",
            "caption": "Chart.js散布図でメーカー別ポジションを可視化（軸選択可能）"
        },
        {
            "title": "5Forces分析",
            "subtitle": "業界構造の5つの競争要因を評価",
            "img": "pb-planner-05-5forces.png",
            "caption": "ダイヤモンド配置で5Forces各要素のスコアを視覚化"
        },
        {
            "title": "価格帯マップ",
            "subtitle": "競合製品の価格分布を一望",
            "img": "pb-planner-09-pricemap-chart.png",
            "caption": "横棒チャートで全メーカーの価格帯を比較、ベース製品位置を表示"
        },
    ]

    for i, s in enumerate(ss_slides):
        add_screenshot_slide(prs, s["title"], s["subtitle"], s["img"], i + 3, s["caption"])

    # P8: PBカード完成
    add_screenshot_slide(
        prs, "PBカード完成", "7項目の基本情報をチャットで設定",
        "pb-planner-10-pbcard.png", 8,
        "アズワン品番・販売価格・JANコード・メーカー型番・入数・キャッチコピー・仕様差分"
    )

    # P9: 仕様差分作成
    add_screenshot_slide(
        prs, "仕様差分作成", "ベース製品からの変更点を記録・管理",
        "pb-planner-12-specdiff.png", 9,
        "仕様の追加・変更・削除を自動ログ。企画書に変更履歴テーブルとして反映"
    )

    # P10: 企画書Word生成
    add_screenshot_slide(
        prs, "企画書(Word)自動生成", "ワンクリックでフレームワーク分析入り企画書を出力",
        "pb-planner-13-word.png", 10,
        "概要 → フレームワーク分析 → PB仕様 → 仕様諸元表（変更履歴付）の4章構成"
    )

    # P11: マスタExcel完成
    add_screenshot_slide(
        prs, "PIMデータ(Excel)同時完成", "商品マスタ登録用データを企画と同時に生成",
        "pb-planner-14-excel.png", 11,
        "品番・JANコード・価格・スペック26項目・キャッチコピー・PB差分仕様を一括出力"
    )

    # 保存
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"[OK] 提案書を保存: {OUTPUT}")
    print(f"     スライド数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
